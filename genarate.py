import pandas as pd
from datetime import datetime
import requests
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import comtypes.client

# Function to extract necessary information from the trucklist
def extract_trucklist_info(trucklist_file):
    try:
        # Read the Excel file
        df = pd.read_excel(trucklist_file)
        # Filter orders where labels have been created
        orders_to_process = df[df['Labels created?'] == 'Yes']
        return orders_to_process
    except Exception as e:
        print(f"Error extracting trucklist info: {e}")
        raise

# Function to generate a barcode
def generate_barcode(order_number, position):
    barcode_data = f"00{order_number}{position}"
    barcode_url = f"https://barcode.tec-it.com/barcode.ashx?data={barcode_data}&code=Code39"
    try:
        # Download the barcode image
        response = requests.get(barcode_url, timeout=100)
        response.raise_for_status()
        # Save the barcode image locally
        with open('barcode.png', 'wb') as f:
            f.write(response.content)
    except requests.RequestException as e:
        print(f"Error fetching barcode: {e}")
        raise

# Function to create a label using a template
def create_label(order_info, template_folder, output_folder):
    sku = order_info['SKU']
    sku_full = sku
    template_subfolder = os.path.join(template_folder, sku[:5])

    # Ensure the subfolder exists
    if not os.path.isdir(template_subfolder):
        raise FileNotFoundError(f"Template subfolder not found: {template_subfolder}")

    # Find the exact template file within the subfolder
    template_file = None
    for file in os.listdir(template_subfolder):
        if file.startswith(f"Labels {sku_full}") and file.endswith(".pptx"):
            template_file = os.path.join(template_subfolder, file)
            break

    if template_file is None:
        raise FileNotFoundError(f"Template file not found for SKU: {sku_full} in {template_subfolder}")

    print(f"Using template: {template_file}")

    # Load the template and make necessary replacements
    prs = Presentation(template_file)
    
    order_id = str(order_info['Ext order number'])
    position = f"00{int(float(order_info['position'])):03d}00"  # Convert to float then int before formatting

    # Convert order date to string if it is a Timestamp
    order_date = order_info['Order date']
    if isinstance(order_date, pd.Timestamp):
        order_date = order_date.strftime('%d.%m.%Y')
    else:
        order_date = datetime.strptime(order_date, '%d.%m.%Y').strftime('%d.%m.%Y')

    # Text replacement with specific formatting
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if 'ORDER_ID' in run.text:
                            run.text = run.text.replace('ORDER_ID', order_id)
                            run.font.size = Pt(7.5)
                            run.font.name = 'Arial'
                        if 'POSITION' in run.text:
                            run.text = run.text.replace('POSITION', position)
                            run.font.size = Pt(7.5)
                            run.font.name = 'Arial'
                        if 'ORDER_DATE' in run.text:
                            run.text = run.text.replace('ORDER_DATE', order_date)
                            run.font.size = Pt(7.5)
                            run.font.name = 'Arial'
                        # Adjust the formatting for other specific text elements
                        if 'Set artikel:' in run.text or 'Artnr.leverencier' in run.text:
                            run.font.size = Pt(9)
                            run.font.name = 'Arial'
                        if 'kg' in run.text or '1/' in run.text:
                            run.font.size = Pt(9)
                            run.font.name = 'Arial'
                        if 'Emma' in run.text or sku_full in run.text:
                            run.font.size = Pt(10)
                            run.font.name = 'Arial'
                        if 'Bestelling/positie:' in run.text or 'Ontvangsdatum:' in run.text:
                            run.font.size = Pt(7.5)
                            run.font.name = 'Arial'

    # Insert the barcode image
    generate_barcode(order_id, position)
    barcode_image_path = 'barcode.png'
    for slide in prs.slides:
        # Adjust the position and size of the barcode to fit the top right box
        slide.shapes.add_picture(barcode_image_path, Inches(5.0), Inches(0.5), width=Inches(2.0), height=Inches(0.75))

    # Create a new subfolder for each new order if it does not exist
    order_folder = os.path.join(output_folder, order_id)
    if not os.path.exists(order_folder):
        os.makedirs(order_folder)
    
    # Save the presentation
    label_filename = os.path.abspath(os.path.join(order_folder, f"{sku}_{position}.pptx"))
    prs.save(label_filename)

    # Convert the PowerPoint to PDF
    pdf_filename = label_filename.replace('.pptx', '.pdf')
    ppt_to_pdf(label_filename, pdf_filename)

    # Delete the intermediate .pptx file
    try:
        os.remove(label_filename)
    except Exception as e:
        print(f"Error deleting intermediate .pptx file: {e}")

    return pdf_filename

# Function to convert pptx to pdf using comtypes
def ppt_to_pdf(input_file, output_file):
    print(f"Converting {input_file} to {output_file}")
    if not os.path.exists(input_file):
        raise FileNotFoundError(f"Input file not found: {input_file}")

    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1
        deck = powerpoint.Presentations.Open(input_file, WithWindow=False)
        deck.SaveAs(output_file, 32)  # 32 is the formatType for ppt to pdf
        deck.Close()
        powerpoint.Quit()
    except Exception as e:
        print(f"Error converting PPT to PDF: {e}")
        raise

# Function to update the trucklist
def update_trucklist(trucklist_file, orders_processed):
    try:
        df = pd.read_excel(trucklist_file)
        df.loc[df['Order number'].isin(orders_processed), 'Labels created?'] = 'No'
        df.to_excel(trucklist_file, index=False)
    except Exception as e:
        print(f"Error updating trucklist: {e}")
        raise

# Main function to orchestrate the workflow
def main():
    trucklist_file = os.path.abspath('TruckList.xlsx')
    template_folder = os.path.abspath('labels')
    output_folder = os.path.abspath('generated_labels')
    
    print(f"Current directory: {os.getcwd()}")
    print(f"Files in trucklist directory: {os.listdir(os.path.dirname(trucklist_file))}")
    print(f"Files in template directory: {os.listdir(template_folder)}")
    
    try:
        orders_to_process = extract_trucklist_info(trucklist_file)
    except Exception as e:
        print(f"Failed to extract trucklist info: {e}")
        return
    
    orders_processed = []
    
    for _, order_info in orders_to_process.iterrows():
        try:
            label_filename = create_label(order_info, template_folder, output_folder)
            print(f"Label created: {label_filename}")
            orders_processed.append(order_info['Order number'])
        except Exception as e:
            print(f"Failed to create label for order {order_info['Order number']}: {e}")

    try:
        update_trucklist(trucklist_file, orders_processed)
    except Exception as e:
        print(f"Failed to update trucklist: {e}")

if __name__ == "__main__":
    main()
